# -*- coding: utf-8 -*-
"""
Canvaè‡ªå‹•åŒ–ãƒãƒ³ãƒ‰ãƒ©ãƒ¼ï¼ˆRailwayç‰ˆï¼‰
WooCommerce Webhookã‚’å—ä¿¡ã—ã¦Canvaãƒ‡ã‚¶ã‚¤ãƒ³ã‚’è‡ªå‹•ä½œæˆ
"""
import os
import re
import json
import base64
import requests
import time
import tempfile
from io import BytesIO
from datetime import datetime

# python-pptx
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Pillow
from PIL import Image

# reportlab for PDF creation
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas as pdf_canvas
from reportlab.lib.utils import ImageReader
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont

# è¨­å®šï¼ˆç’°å¢ƒå¤‰æ•°ã‹ã‚‰å–å¾— - é…å»¶èª­ã¿è¾¼ã¿ï¼‰
def get_canva_client_id():
    return os.getenv("CANVA_CLIENT_ID", "OC-AZvUVtxGhbOD")

def get_canva_client_secret():
    return os.getenv("CANVA_CLIENT_SECRET", "")

# ã‚µãƒ¼ãƒãƒ¼ä¸Šã®cutoutç”»åƒURL
CUTOUT_BASE_URL = "https://i-tategu-shop.com/wp-content/themes/i-tategu/assets/images/cutouts"
TREE_IMAGES_URL = "https://i-tategu-shop.com/wp-content/themes/i-tategu/assets/images"

# ã‚¹ãƒ©ã‚¤ãƒ‰ã‚µã‚¤ã‚º
SLIDE_WIDTH_PX = 1000
SLIDE_HEIGHT_PX = 1000
EMU_PER_PX = 914400 / 96

# ãƒ•ã‚©ãƒ³ãƒˆãƒãƒƒãƒ”ãƒ³ã‚°
FONT_MAP = {
    'Alex Brush': 'Alex Brush',
    'Great Vibes': 'Great Vibes',
    'Pinyon Script': 'Pinyon Script',
    'Sacramento': 'Sacramento',
    'Dancing Script': 'Dancing Script',
    'Parisienne': 'Parisienne',
    'Allura': 'Allura',
    'Satisfy': 'Satisfy',
    'Rouge Script': 'Rouge Script',
    # æ—¥æœ¬èªãƒ•ã‚©ãƒ³ãƒˆï¼ˆäººæ°—ï¼‰
    'Shippori Mincho': 'Shippori Mincho',
    'Zen Old Mincho': 'Zen Old Mincho',
    'Klee One': 'Klee One',
    'Noto Serif JP': 'Noto Serif JP',
    'Zen Maru Gothic': 'Zen Maru Gothic',
    'Sawarabi Mincho': 'Sawarabi Mincho',
    'Noto Sans JP': 'Noto Sans JP',
    # æ—¥æœ¬èªãƒ•ã‚©ãƒ³ãƒˆï¼ˆå€‹æ€§æ´¾ï¼‰
    'Yomogi': 'Yomogi',
    'Kaisei Decol': 'Kaisei Decol',
    'Reggae One': 'Reggae One',
}

# ãƒ•ã‚©ãƒ³ãƒˆè¡¨ç¤ºåãƒãƒƒãƒ”ãƒ³ã‚°ï¼ˆã‚·ãƒŸãƒ¥ãƒ¬ãƒ¼ã‚¿ãƒ¼ã¨åŒã˜ã‚¹ã‚¿ã‚¤ãƒ«è¡¨è¨˜ã€è‹±èªç‰ˆï¼‰
FONT_DISPLAY_MAP = {
    'Sacramento': 'Holiday style',
    'Pinyon Script': 'Eyesome style',
    'Satisfy': 'Mistrully style',
    'Rouge Script': 'Amsterdam style',
}

# ãƒ†ãƒ³ãƒ—ãƒ¬ãƒ¼ãƒˆ
TEMPLATES = {
    'holy': """I hereby certify
that this man and woman
were united in holy matrimony,
in the name of the Father,
of the Son and of the Holy Spirit.""",
    'happy': """We are happy to be married
before Holy God and many witnesses.
We promise that in this new life
we will have joy and happiness.""",
    'promise': """We promise to be husband and wife
before these witnesses here present.
We swear to admire, love and help each other
and to create a peaceful happy family."""
}

TITLES = {
    'wedding': 'Wedding Certificate',
    'marriage': 'Certificate of Marriage'
}

BACKGROUND_MAP = {
    'product': '_1_cutout.png',
    'product_back': '_2_cutout.png',
    'noclear': '_3_cutout.png',
    'noclear_back': '_4_cutout.png',
}


def px_to_emu(px):
    return int(px * EMU_PER_PX)


def format_date(date_str, format_type):
    try:
        for fmt in ["%Y-%m-%d", "%Y.%m.%d", "%Y/%m/%d"]:
            try:
                d = datetime.strptime(date_str, fmt)
                break
            except:
                continue
        else:
            return date_str
    except:
        return date_str

    months_full = ['January', 'February', 'March', 'April', 'May', 'June',
                   'July', 'August', 'September', 'October', 'November', 'December']
    months_short = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun',
                    'Jul', 'Aug', 'Sep', 'Oct', 'Nov', 'Dec']

    def ordinal(n):
        s = ['th', 'st', 'nd', 'rd']
        v = n % 100
        return str(n) + (s[(v - 20) % 10] if v > 20 else s[v] if v < 4 else s[0])

    if format_type == 'western':
        return f"{d.year}.{d.month:02d}.{d.day:02d}"
    elif format_type == 'us_long':
        return f"{months_full[d.month-1]} {ordinal(d.day)}, {d.year}"
    elif format_type == 'us_short':
        return f"{months_short[d.month-1]} {ordinal(d.day)}, {d.year}"
    elif format_type == 'uk_long':
        return f"{ordinal(d.day)} {months_full[d.month-1]} {d.year}"
    elif format_type == 'uk_short':
        return f"{ordinal(d.day)} {months_short[d.month-1]} {d.year}"
    else:
        return date_str


def add_text_box(slide, text, x_px, y_px, font_name, font_size_pt, center=True, color_rgb=(42, 24, 16)):
    width_px = max(len(text) * font_size_pt * 0.8, 100)
    height_px = font_size_pt * 2

    left = px_to_emu(x_px - width_px / 2) if center else px_to_emu(x_px)
    top = px_to_emu(y_px - height_px / 2)
    width = px_to_emu(width_px)
    height = px_to_emu(height_px)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = False

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER if center else PP_ALIGN.LEFT

    run = p.add_run()
    run.text = text
    run.font.name = FONT_MAP.get(font_name, font_name)
    run.font.size = Pt(font_size_pt)
    run.font.color.rgb = RGBColor(*color_rgb)

    return textbox


def add_multiline_text_box(slide, text, x_px, y_px, font_name, font_size_pt, line_height=1.4, color_rgb=(42, 24, 16)):
    lines = text.strip().split('\n')
    width_px = max(len(line) * font_size_pt * 0.7 for line in lines)
    height_px = len(lines) * font_size_pt * line_height * 1.5

    left = px_to_emu(x_px - width_px / 2)
    top = px_to_emu(y_px)
    width = px_to_emu(width_px)
    height = px_to_emu(height_px)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = False

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(0)
        p.space_after = Pt(font_size_pt * (line_height - 1))

        run = p.add_run()
        run.text = line.strip()
        run.font.name = FONT_MAP.get(font_name, font_name)
        run.font.size = Pt(font_size_pt)
        run.font.color.rgb = RGBColor(*color_rgb)

    return textbox


def extract_board_info(product_name):
    """å•†å“åã‹ã‚‰æ¿æƒ…å ±ã‚’æŠ½å‡º"""
    name = ''
    number = '01'
    size = ''

    # ãƒ‘ã‚¿ãƒ¼ãƒ³1: "ã€ä¸€ç‚¹ç‰©ã€‘ã‚¿ãƒ¢ ä¸€æšæ¿ çµå©šè¨¼æ˜æ›¸ 300x300mm"
    match = re.search(r'ã€ä¸€ç‚¹ç‰©ã€‘\s*(.+?)\s+ä¸€æšæ¿.*?(\d+)x(\d+)', product_name)
    if match:
        name = match.group(1).strip()
        size = f"{match.group(2)}_{match.group(3)}"
        return {'name': name, 'number': number, 'size': size}

    # ãƒ‘ã‚¿ãƒ¼ãƒ³2: "ã‚±ãƒ¤ã‚­_01_400_600"
    match = re.match(r'^([^_]+)_(\d+)(?:_(\d+)_(\d+))?', product_name)
    if match:
        name = match.group(1)
        number = match.group(2)
        if match.group(3) and match.group(4):
            size = f"{match.group(3)}_{match.group(4)}"
        return {'name': name, 'number': number.zfill(2), 'size': size}

    # ãƒ‘ã‚¿ãƒ¼ãƒ³3: "ã‚±ãƒ¤ã‚­ No.01"
    match = re.search(r'^(.+?)\s*No\.?(\d+)', product_name)
    if match:
        name = match.group(1).strip()
        number = match.group(2)
        return {'name': name, 'number': number.zfill(2), 'size': size}

    return {'name': name, 'number': number, 'size': size}


def find_cutout_url(board_name, board_number, board_size, background='product'):
    """cutoutç”»åƒã®URLã‚’æ§‹ç¯‰"""
    suffix = BACKGROUND_MAP.get(background, '_1_cutout.png')

    # ã‚µã‚¤ã‚ºã‚ã‚Šã®å ´åˆ
    if board_size:
        filename = f"{board_name}_{board_number}_{board_size}{suffix}"
    else:
        filename = f"{board_name}_{board_number}_300_300{suffix}"  # ãƒ‡ãƒ•ã‚©ãƒ«ãƒˆã‚µã‚¤ã‚º

    return f"{CUTOUT_BASE_URL}/{filename}"


def download_image(url, temp_dir, max_size=800, preserve_transparency=False):
    """URLã‹ã‚‰ç”»åƒã‚’ãƒ€ã‚¦ãƒ³ãƒ­ãƒ¼ãƒ‰ã—ã¦åœ§ç¸®ä¿å­˜

    Args:
        preserve_transparency: Trueã®å ´åˆã€PNGå½¢å¼ã§é€æ˜åº¦ã‚’ä¿æŒ
    """
    try:
        response = requests.get(url, timeout=30)
        if response.status_code == 200:
            # ç”»åƒã‚’é–‹ã„ã¦ãƒªã‚µã‚¤ã‚ºãƒ»åœ§ç¸®
            img = Image.open(BytesIO(response.content))

            # ãƒªã‚µã‚¤ã‚ºï¼ˆæœ€å¤§ã‚µã‚¤ã‚ºã‚’è¶…ãˆã‚‹å ´åˆï¼‰
            if max(img.size) > max_size:
                ratio = max_size / max(img.size)
                new_size = (int(img.size[0] * ratio), int(img.size[1] * ratio))
                img = img.resize(new_size, Image.LANCZOS)

            if preserve_transparency:
                # PNGå½¢å¼ã§é€æ˜åº¦ã‚’ä¿æŒ
                if img.mode != 'RGBA':
                    img = img.convert('RGBA')
                filename = os.path.splitext(os.path.basename(url))[0] + '.png'
                filepath = os.path.join(temp_dir, filename)
                img.save(filepath, 'PNG', optimize=True)
            else:
                # JPEGå½¢å¼ï¼ˆé€æ˜åº¦ã‚’ç™½èƒŒæ™¯ã«å¤‰æ›ï¼‰
                if img.mode == 'RGBA':
                    white_bg = Image.new('RGB', img.size, (255, 255, 255))
                    white_bg.paste(img, mask=img.split()[3])
                    img = white_bg
                elif img.mode != 'RGB':
                    img = img.convert('RGB')
                filename = os.path.splitext(os.path.basename(url))[0] + '.jpg'
                filepath = os.path.join(temp_dir, filename)
                img.save(filepath, 'JPEG', quality=80, optimize=True)

            print(f"[IMG] {'PNG' if preserve_transparency else 'JPEG'}: {os.path.getsize(filepath) / 1024:.1f}KB")
            return filepath
    except Exception as e:
        print(f"[WARN] Failed to download image: {url} - {e}")
    return None


def refresh_canva_token(refresh_token):
    """Canvaãƒˆãƒ¼ã‚¯ãƒ³ã‚’ãƒªãƒ•ãƒ¬ãƒƒã‚·ãƒ¥ã—ã€ç’°å¢ƒå¤‰æ•°ã«ä¿å­˜"""
    url = 'https://api.canva.com/rest/v1/oauth/token'
    client_id = get_canva_client_id()
    client_secret = get_canva_client_secret()

    print(f"[Canva Token] Client ID: {client_id}")
    print(f"[Canva Token] Client Secret: {'SET' if client_secret else 'EMPTY'} (len={len(client_secret) if client_secret else 0})")

    credentials = base64.b64encode(f"{client_id}:{client_secret}".encode()).decode()

    response = requests.post(url, data={
        'grant_type': 'refresh_token',
        'refresh_token': refresh_token,
    }, headers={
        'Authorization': f'Basic {credentials}',
        'Content-Type': 'application/x-www-form-urlencoded',
    })

    print(f"[Canva Token] Refresh status: {response.status_code}")

    if response.status_code == 200:
        tokens = response.json()
        new_access = tokens.get('access_token')
        new_refresh = tokens.get('refresh_token', refresh_token)

        # ç’°å¢ƒå¤‰æ•°ã‚’æ›´æ–°ï¼ˆãƒ—ãƒ­ã‚»ã‚¹å†…ã§æ°¸ç¶šåŒ–ï¼‰
        os.environ['CANVA_ACCESS_TOKEN'] = new_access
        os.environ['CANVA_REFRESH_TOKEN'] = new_refresh
        print(f"[Canva Token] Refresh successful! Tokens updated in memory.")
        print(f"[Canva Token] New refresh token (first 50 chars): {new_refresh[:50]}...")

        return {
            'access_token': new_access,
            'refresh_token': new_refresh
        }

    print(f"[Canva Token] Refresh failed: {response.text[:500]}")
    return None


def import_to_canva(file_path, title, access_token, refresh_token, retry=False):
    """Canvaã«ãƒ•ã‚¡ã‚¤ãƒ«ã‚’ã‚¤ãƒ³ãƒãƒ¼ãƒˆPDF/PPTXå¯¾å¿œï¼ˆã‚¨ãƒ©ãƒ¼æƒ…å ±ã‚‚è¿”ã™ï¼‰"""
    url = "https://api.canva.com/rest/v1/imports"
    title_base64 = base64.b64encode(title.encode("utf-8")).decode("utf-8")

    # ãƒ•ã‚¡ã‚¤ãƒ«å½¢å¼ã‚’æ¤œå‡ºã—ã¦MIMEã‚¿ã‚¤ãƒ—ã‚’è¨­å®š
    if file_path.endswith('.pdf'):
        mime_type = "application/pdf"
        print(f"[Canva Import] Using PDF format")
    else:
        mime_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        print(f"[Canva Import] Using PPTX format")

    metadata = {
        "title_base64": title_base64,
        "mime_type": mime_type
    }

    headers = {
        "Authorization": f"Bearer {access_token}",
        "Content-Type": "application/octet-stream",
        "Import-Metadata": json.dumps(metadata)
    }

    with open(file_path, "rb") as f:
        response = requests.post(url, headers=headers, data=f)

    print(f"[Canva Import] Status: {response.status_code}")

    # 401ã‚¨ãƒ©ãƒ¼æ™‚ã¯ãƒˆãƒ¼ã‚¯ãƒ³ãƒªãƒ•ãƒ¬ãƒƒã‚·ãƒ¥
    if response.status_code == 401 and not retry:
        print("[INFO] Token expired, refreshing...")
        new_tokens = refresh_canva_token(refresh_token)
        if new_tokens:
            return import_to_canva(file_path, title, new_tokens['access_token'], new_tokens['refresh_token'], retry=True)
        return None, {"error": "Token refresh failed"}

    if response.status_code != 200:
        error_msg = f"HTTP {response.status_code}: {response.text[:500]}"
        print(f"[ERROR] Canva import failed: {error_msg}")
        return None, {"error": error_msg}

    job = response.json().get("job", {})
    job_id = job.get("id")
    print(f"[Canva Import] Job ID: {job_id}")

    # ãƒªãƒ•ãƒ¬ãƒƒã‚·ãƒ¥ã—ãŸå ´åˆã¯æ–°ã—ã„ãƒˆãƒ¼ã‚¯ãƒ³ã‚’ä½¿ç”¨
    check_token = access_token

    # ã‚¸ãƒ§ãƒ–å®Œäº†ã‚’å¾…æ©Ÿ
    for i in range(15):
        time.sleep(2)
        check_url = f"https://api.canva.com/rest/v1/imports/{job_id}"
        check_resp = requests.get(check_url, headers={"Authorization": f"Bearer {check_token}"})

        print(f"[Canva Import] Check {i+1}: {check_resp.status_code}")

        if check_resp.status_code == 200:
            job_data = check_resp.json().get("job", {})
            status = job_data.get("status")
            print(f"[Canva Import] Status: {status}")

            if status == "success":
                designs = job_data.get("result", {}).get("designs", [])
                if designs:
                    return designs[0], None
            elif status == "failed":
                error = job_data.get("error", {})
                error_msg = error.get('message', 'Unknown error')
                print(f"[ERROR] {error_msg}")
                return None, {"error": error_msg, "details": error}
        elif check_resp.status_code == 401:
            # ãƒã‚§ãƒƒã‚¯ä¸­ã«ãƒˆãƒ¼ã‚¯ãƒ³ãŒåˆ‡ã‚ŒãŸå ´åˆã€ãƒªãƒ•ãƒ¬ãƒƒã‚·ãƒ¥
            new_tokens = refresh_canva_token(refresh_token)
            if new_tokens:
                check_token = new_tokens['access_token']

    print("[ERROR] Timeout waiting for Canva import")
    return None, {"error": "Timeout"}


def get_order_from_woocommerce(order_id, wc_url, wc_key, wc_secret):
    """WooCommerceã‹ã‚‰æ³¨æ–‡ã‚’å–å¾—"""
    url = f"{wc_url}/wp-json/wc/v3/orders/{order_id}"
    print(f"[WC API] Fetching: {url}")

    try:
        response = requests.get(url, auth=(wc_key, wc_secret))
        print(f"[WC API] Status: {response.status_code}")

        if response.status_code == 200:
            return response.json()
        else:
            print(f"[WC API] Error response: {response.text[:500]}")
            return None
    except Exception as e:
        print(f"[WC API] Exception: {e}")
        return None


def parse_order_data(order):
    """æ³¨æ–‡ãƒ‡ãƒ¼ã‚¿ã‚’è§£æ"""
    meta = {m['key']: m['value'] for m in order.get('meta_data', [])}

    # ã‚·ãƒŸãƒ¥ãƒ¬ãƒ¼ã‚·ãƒ§ãƒ³ãƒ‡ãƒ¼ã‚¿
    sim_data_raw = meta.get('_simulation_data', '{}')
    try:
        sim_data = json.loads(sim_data_raw) if sim_data_raw else {}
    except:
        sim_data = {}

    # ã‚·ãƒŸãƒ¥ãƒ¬ãƒ¼ã‚·ãƒ§ãƒ³ç”»åƒï¼ˆé€ã‹ã—ãªã—ãƒ•ã‚¡ã‚¤ãƒ«URLã‚’å„ªå…ˆï¼‰
    sim_image_url = meta.get('_simulation_image_url', '')
    sim_image = sim_image_url if sim_image_url else meta.get('_simulation_image', '')

    # åˆ»å°æƒ…å ±ï¼ˆãƒ•ã‚©ãƒ¼ãƒ«ãƒãƒƒã‚¯ï¼‰
    if not sim_data.get('groomName'):
        sim_data['groomName'] = meta.get('_engraving_name1', '')
    if not sim_data.get('brideName'):
        sim_data['brideName'] = meta.get('_engraving_name2', '')

    # ãƒ•ã‚©ãƒ³ãƒˆæƒ…å ±ï¼ˆãƒ•ã‚©ãƒ¼ãƒ«ãƒãƒƒã‚¯ï¼‰
    if not sim_data.get('baseFont'):
        font_style = meta.get('_font_style', '')
        if font_style:
            sim_data['baseFont'] = font_style

    # å•†å“æƒ…å ±
    product_name = ''
    for item in order.get('line_items', []):
        product_name = item.get('name', '')
        break

    board_info = extract_board_info(product_name)

    # æ—¥ä»˜
    wedding_date = meta.get('_engraving_date', '')
    if not wedding_date:
        wedding_date = sim_data.get('weddingDate', '')

    return {
        'order_id': order['id'],
        'sim_data': sim_data,
        'sim_image': sim_image,
        'product_name': product_name,
        'board_name': board_info['name'],
        'board_number': board_info['number'],
        'board_size': board_info['size'],
        'wedding_date': wedding_date,
    }


def create_pptx(order_data, temp_dir):
    """PowerPointã‚’ä½œæˆï¼ˆ6ãƒšãƒ¼ã‚¸æ§‹æˆï¼‰"""
    print(f"[Canva] Creating PowerPoint for order #{order_data['order_id']}...")

    sim_data = order_data['sim_data']
    sim_image = order_data['sim_image']
    groom = sim_data.get('groomName', '')
    bride = sim_data.get('brideName', '')

    prs = Presentation()
    prs.slide_width = px_to_emu(SLIDE_WIDTH_PX)
    prs.slide_height = px_to_emu(SLIDE_HEIGHT_PX)
    blank_layout = prs.slide_layouts[6]

    # cutoutç”»åƒã®URLã‚’å–å¾—
    background = sim_data.get('background', 'product')
    cutout_urls = {
        'product': find_cutout_url(order_data['board_name'], order_data['board_number'], order_data['board_size'], 'product'),
        'product_back': find_cutout_url(order_data['board_name'], order_data['board_number'], order_data['board_size'], 'product_back'),
        'noclear': find_cutout_url(order_data['board_name'], order_data['board_number'], order_data['board_size'], 'noclear'),
        'noclear_back': find_cutout_url(order_data['board_name'], order_data['board_number'], order_data['board_size'], 'noclear_back'),
    }

    # ========== 1ãƒšãƒ¼ã‚¸ç›®: ã‚·ãƒŸãƒ¥ãƒ¬ãƒ¼ã‚·ãƒ§ãƒ³ç”»åƒ + æ³¨æ–‡æƒ…å ± ==========
    slide1 = prs.slides.add_slide(blank_layout)

    if sim_image:
        try:
            temp_img_path = os.path.join(temp_dir, f"sim_{order_data['order_id']}.jpg")

            # ç”»åƒãƒ‡ãƒ¼ã‚¿ã‚’å–å¾—
            if sim_image.startswith('http'):
                response = requests.get(sim_image, timeout=30)
                response.raise_for_status()
                img = Image.open(BytesIO(response.content))
            else:
                if ',' in sim_image:
                    sim_image = sim_image.split(',')[1]
                img_data = base64.b64decode(sim_image)
                img = Image.open(BytesIO(img_data))

            # åœ§ç¸®: RGBAâ†’RGBå¤‰æ›ã€ãƒªã‚µã‚¤ã‚º
            if img.mode == 'RGBA':
                white_bg = Image.new('RGB', img.size, (255, 255, 255))
                white_bg.paste(img, mask=img.split()[3])
                img = white_bg
            elif img.mode != 'RGB':
                img = img.convert('RGB')

            # æœ€å¤§800pxã«ãƒªã‚µã‚¤ã‚º
            max_size = 800
            if max(img.size) > max_size:
                ratio = max_size / max(img.size)
                new_size = (int(img.size[0] * ratio), int(img.size[1] * ratio))
                img = img.resize(new_size, Image.LANCZOS)

            # JPEGã§ä¿å­˜
            img.save(temp_img_path, 'JPEG', quality=80, optimize=True)
            print(f"[IMG] Sim image: {os.path.getsize(temp_img_path) / 1024:.1f}KB")

            img = Image.open(temp_img_path)
            img_width, img_height = img.size
            available_height = SLIDE_HEIGHT_PX * 0.75
            scale = min(SLIDE_WIDTH_PX * 0.95 / img_width, available_height / img_height)
            draw_width = img_width * scale
            draw_height = img_height * scale
            x = (SLIDE_WIDTH_PX - draw_width) / 2
            y = 20

            slide1.shapes.add_picture(
                temp_img_path,
                px_to_emu(x), px_to_emu(y),
                px_to_emu(draw_width), px_to_emu(draw_height)
            )
        except Exception as e:
            print(f"[WARN] Simulation image error: {e}")

    # ä¸‹éƒ¨ã«æ³¨æ–‡æƒ…å ±
    info_y = SLIDE_HEIGHT_PX - 160
    add_text_box(slide1, f"æ³¨æ–‡ #{order_data['order_id']} - {order_data['board_name']} No.{order_data['board_number']}",
                 SLIDE_WIDTH_PX/2, info_y, 'Meiryo UI', 16, center=True, color_rgb=(40, 40, 40))

    info_y += 30
    add_text_box(slide1, f"æ–°éƒ: {groom}ã€€ã€€æ–°å©¦: {bride}",
                 SLIDE_WIDTH_PX/2, info_y, 'Meiryo UI', 13, center=True, color_rgb=(60, 60, 60))

    info_y += 26
    add_text_box(slide1, f"æŒ™å¼æ—¥: {order_data['wedding_date']}",
                 SLIDE_WIDTH_PX/2, info_y, 'Meiryo UI', 13, center=True, color_rgb=(60, 60, 60))

    info_y += 26
    base_font_display = sim_data.get('baseFont', 'Alex Brush')
    template_names = {'holy': 'æ•™ä¼šå¼â‘ ', 'happy': 'æ•™ä¼šå¼â‘¡', 'promise': 'äººå‰å¼', 'custom': 'ã‚«ã‚¹ã‚¿ãƒ '}
    template_info = template_names.get(sim_data.get('template', 'holy'), 'æ•™ä¼šå¼â‘ ')

    # å„è¦ç´ ã®ãƒ•ã‚©ãƒ³ãƒˆã‚’ãƒã‚§ãƒƒã‚¯ï¼ˆãƒ™ãƒ¼ã‚¹ã¨ç•°ãªã‚‹å ´åˆã®ã¿è¡¨ç¤ºï¼‰
    font_differences = []
    title_font = sim_data.get('titleFont', '')
    body_font = sim_data.get('bodyFont', '')
    date_font = sim_data.get('dateFont', '')
    name_font = sim_data.get('nameFont', '')

    if title_font and title_font != base_font_display:
        font_differences.append(f"ã‚¿ã‚¤ãƒˆãƒ«:{title_font}")
    if body_font and body_font != base_font_display:
        font_differences.append(f"æœ¬æ–‡:{body_font}")
    if date_font and date_font != base_font_display:
        font_differences.append(f"æ—¥ä»˜:{date_font}")
    if name_font and name_font != base_font_display:
        font_differences.append(f"åå‰:{name_font}")

    # ãƒ•ã‚©ãƒ³ãƒˆè¡¨ç¤ºæ–‡å­—åˆ—ã‚’æ§‹ç¯‰
    if font_differences:
        font_detail = f"ãƒ•ã‚©ãƒ³ãƒˆ: {base_font_display}ï¼ˆ{' / '.join(font_differences)}ï¼‰"
    else:
        font_detail = f"ãƒ•ã‚©ãƒ³ãƒˆ: {base_font_display}"

    add_text_box(slide1, f"{font_detail}ã€€ã€€æœ¬æ–‡: {template_info}",
                 SLIDE_WIDTH_PX/2, info_y, 'Meiryo UI', 11, center=True, color_rgb=(100, 100, 100))

    # ========== 2ãƒšãƒ¼ã‚¸ç›®: èƒŒæ™¯cutout + ãƒ†ã‚­ã‚¹ãƒˆ ==========
    slide2 = prs.slides.add_slide(blank_layout)

    base_font = sim_data.get('baseFont', 'Alex Brush')
    if base_font not in FONT_MAP:
        base_font = 'Alex Brush'

    def get_element_font(element):
        font_key = element + 'Font'
        font = sim_data.get(font_key) or base_font
        if font not in FONT_MAP:
            font = base_font
        return font

    FONT_SCALE = SLIDE_WIDTH_PX / 500
    board_size_pct = sim_data.get('boardSize', 130) / 100
    # ãƒ†ã‚­ã‚¹ãƒˆè‰²: 'burn'=èŒ¶è‰², 'white'=ç™½
    text_color_name = sim_data.get('textColor', 'burn')
    text_color = (255, 255, 255) if text_color_name == 'white' else (42, 24, 16)

    # èƒŒæ™¯ç”»åƒã‚’ãƒ€ã‚¦ãƒ³ãƒ­ãƒ¼ãƒ‰ï¼ˆé€æ˜åº¦ã‚’ä¿æŒï¼‰
    cutout_path = download_image(cutout_urls.get(background, cutout_urls['product']), temp_dir, preserve_transparency=True)

    if cutout_path and os.path.exists(cutout_path):
        try:
            bg_img = Image.open(cutout_path)
            bg_width, bg_height = bg_img.size
            base_scale = 0.95
            img_ratio = bg_width / bg_height
            max_width = SLIDE_WIDTH_PX * base_scale
            max_height = SLIDE_HEIGHT_PX * base_scale

            if bg_width / max_width > bg_height / max_height:
                base_width = max_width
                base_height = max_width / img_ratio
            else:
                base_height = max_height
                base_width = max_height * img_ratio

            draw_width = base_width * board_size_pct
            draw_height = base_height * board_size_pct
            img_x = (SLIDE_WIDTH_PX - draw_width) / 2
            img_y = (SLIDE_HEIGHT_PX - draw_height) / 2

            slide2.shapes.add_picture(
                cutout_path,
                px_to_emu(img_x), px_to_emu(img_y),
                px_to_emu(draw_width), px_to_emu(draw_height)
            )
        except Exception as e:
            print(f"[WARN] Background image error: {e}")

    # ã‚¿ã‚¤ãƒˆãƒ«ï¼ˆç‹¬ç«‹ã—ãŸãƒ†ã‚­ã‚¹ãƒˆãƒœãƒƒã‚¯ã‚¹ï¼‰
    title_font = get_element_font('title')
    title_key = sim_data.get('title', 'wedding')
    title_text = sim_data.get('customTitle', '') if title_key == 'custom' else TITLES.get(title_key, 'Wedding Certificate')
    title_x = SLIDE_WIDTH_PX * (sim_data.get('titleX', 50) / 100)
    title_y = SLIDE_HEIGHT_PX * (sim_data.get('titleY', 22) / 100)
    title_size = 24 * (sim_data.get('titleSize', 100) / 100) * FONT_SCALE
    title_box = add_text_box(slide2, title_text, title_x, title_y, title_font, title_size, center=True, color_rgb=text_color)
    # ã‚¿ã‚¤ãƒˆãƒ«ãƒœãƒƒã‚¯ã‚¹ã®ä¸‹ç«¯ã‚’è¨ˆç®—
    title_bottom = title_y + title_size

    # æœ¬æ–‡ï¼ˆç‹¬ç«‹ã—ãŸãƒ†ã‚­ã‚¹ãƒˆãƒœãƒƒã‚¯ã‚¹ - ã‚¿ã‚¤ãƒˆãƒ«ã¨é‡ãªã‚‰ãªã„ã‚ˆã†ã«é…ç½®ï¼‰
    body_font = get_element_font('body')
    template_key = sim_data.get('template', 'holy')
    body_text = sim_data.get('customText', '') if template_key == 'custom' else TEMPLATES.get(template_key, '')
    body_x = SLIDE_WIDTH_PX * (sim_data.get('bodyX', 50) / 100)
    body_y_base = SLIDE_HEIGHT_PX * (sim_data.get('bodyY', 32) / 100)
    # ã‚¿ã‚¤ãƒˆãƒ«ã¨æœ¬æ–‡ãŒé‡ãªã‚‰ãªã„ã‚ˆã†ã«æœ€å°é–“éš”ã‚’ç¢ºä¿
    min_gap = 30  # æœ€å°30pxã®é–“éš”
    body_y = max(body_y_base, title_bottom + min_gap)
    body_size = 11 * (sim_data.get('bodySize', 115) / 100) * FONT_SCALE
    body_line_height = sim_data.get('bodyLineHeight', 1.4)
    if body_text:
        add_multiline_text_box(slide2, body_text, body_x, body_y, body_font, body_size,
                               line_height=body_line_height, color_rgb=text_color)

    # æ—¥ä»˜
    date_font = get_element_font('date')
    date_format_key = sim_data.get('dateFormat', 'western')
    formatted_date = sim_data.get('customDate', '') if date_format_key == 'custom' else format_date(order_data['wedding_date'], date_format_key)
    date_x = SLIDE_WIDTH_PX * (sim_data.get('dateX', 50) / 100)
    date_y = SLIDE_HEIGHT_PX * (sim_data.get('dateY', 60) / 100)
    date_size = 18 * (sim_data.get('dateSize', 85) / 100) * FONT_SCALE
    if formatted_date:
        add_text_box(slide2, formatted_date, date_x, date_y, date_font, date_size, center=True, color_rgb=text_color)

    # åå‰
    name_font = get_element_font('name')
    name_x_pct = sim_data.get('nameX', 50) / 100
    name_y_pct = sim_data.get('nameY', 74) / 100
    name_size = 32 * (sim_data.get('nameSize', 90) / 100) * FONT_SCALE
    name_center_x = SLIDE_WIDTH_PX * name_x_pct
    name_y = SLIDE_HEIGHT_PX * name_y_pct

    groom_width_approx = len(groom) * name_size * 0.6
    amp_width_approx = name_size * 2
    bride_width_approx = len(bride) * name_size * 0.6
    total_width_approx = groom_width_approx + amp_width_approx + bride_width_approx

    if groom:
        groom_x = name_center_x - total_width_approx / 2 + groom_width_approx / 2
        add_text_box(slide2, groom, groom_x, name_y, name_font, name_size, center=True, color_rgb=text_color)

    amp_x = name_center_x
    add_text_box(slide2, "&", amp_x, name_y, name_font, name_size, center=True, color_rgb=text_color)

    if bride:
        bride_x = name_center_x + total_width_approx / 2 - bride_width_approx / 2
        add_text_box(slide2, bride, bride_x, name_y, name_font, name_size, center=True, color_rgb=text_color)

    # ãƒ„ãƒªãƒ¼
    if sim_data.get('showTree', False):
        tree_type = sim_data.get('treeType', 'simple')
        tree_x_pct = sim_data.get('treeX', 0.75)
        tree_y_pct = sim_data.get('treeY', 0.65)
        tree_size_pct = sim_data.get('treeSize', 80) / 100

        tree_url = f"{TREE_IMAGES_URL}/tree-{tree_type}.png"
        print(f"[TREE] Downloading: {tree_url}")

        try:
            # ç›´æ¥ãƒ€ã‚¦ãƒ³ãƒ­ãƒ¼ãƒ‰ã—ã¦PNGå½¢å¼ã‚’æ˜ç¤ºçš„ã«ä¿æŒ
            response = requests.get(tree_url, timeout=30)
            if response.status_code == 200:
                tree_img = Image.open(BytesIO(response.content))
                print(f"[TREE] Original: {tree_img.size}, mode={tree_img.mode}")

                # ãƒªã‚µã‚¤ã‚º
                max_size = 800
                if max(tree_img.size) > max_size:
                    ratio = max_size / max(tree_img.size)
                    new_size = (int(tree_img.size[0] * ratio), int(tree_img.size[1] * ratio))
                    tree_img = tree_img.resize(new_size, Image.LANCZOS)

                # RGBAç¢ºä¿
                if tree_img.mode != 'RGBA':
                    tree_img = tree_img.convert('RGBA')
                print(f"[TREE] After resize: {tree_img.size}, mode={tree_img.mode}")

                # é€æ˜åº¦ãƒã‚§ãƒƒã‚¯
                alpha = tree_img.split()[3]
                alpha_extrema = alpha.getextrema()
                print(f"[TREE] Alpha range: {alpha_extrema}")

                # BytesIOã«PNGã¨ã—ã¦ä¿å­˜
                tree_buffer = BytesIO()
                tree_img.save(tree_buffer, 'PNG', optimize=False)
                tree_buffer.seek(0)
                print(f"[TREE] PNG buffer size: {len(tree_buffer.getvalue()) / 1024:.1f}KB")

                tree_width, tree_height = tree_img.size
                base_tree_size = min(SLIDE_WIDTH_PX, SLIDE_HEIGHT_PX) * 0.3
                draw_tree_width = base_tree_size * tree_size_pct
                draw_tree_height = draw_tree_width * (tree_height / tree_width)
                tree_x = SLIDE_WIDTH_PX * tree_x_pct - draw_tree_width / 2
                tree_y = SLIDE_HEIGHT_PX * tree_y_pct - draw_tree_height / 2

                # BytesIOã‹ã‚‰ç›´æ¥è¿½åŠ 
                slide2.shapes.add_picture(
                    tree_buffer,
                    px_to_emu(tree_x), px_to_emu(tree_y),
                    px_to_emu(draw_tree_width), px_to_emu(draw_tree_height)
                )
                print(f"[TREE] Added to slide at ({tree_x:.0f}, {tree_y:.0f})")
        except Exception as e:
            print(f"[WARN] Tree image error: {e}")
            import traceback
            traceback.print_exc()

    # ========== 3-6ãƒšãƒ¼ã‚¸ç›®: cutoutç”»åƒ ==========
    cutout_labels = [
        ('product', 'æ°´å¼•è¡¨'),
        ('product_back', 'æ°´å¼•è£'),
        ('noclear', 'ç„¡å¡—è£…è¡¨'),
        ('noclear_back', 'ç„¡å¡—è£…è£'),
    ]

    for key, label in cutout_labels:
        slide = prs.slides.add_slide(blank_layout)
        cutout_path = download_image(cutout_urls[key], temp_dir, preserve_transparency=True)

        if cutout_path and os.path.exists(cutout_path):
            try:
                img = Image.open(cutout_path)
                img_width, img_height = img.size
                max_size = min(SLIDE_WIDTH_PX, SLIDE_HEIGHT_PX) * 0.9
                scale = min(max_size / img_width, max_size / img_height)
                draw_width = img_width * scale
                draw_height = img_height * scale
                x = (SLIDE_WIDTH_PX - draw_width) / 2
                y = (SLIDE_HEIGHT_PX - draw_height) / 2

                slide.shapes.add_picture(
                    cutout_path,
                    px_to_emu(x), px_to_emu(y),
                    px_to_emu(draw_width), px_to_emu(draw_height)
                )
            except Exception as e:
                print(f"[WARN] {label} image error: {e}")
                add_text_box(slide, f"No Image: {label}", SLIDE_WIDTH_PX/2, SLIDE_HEIGHT_PX/2,
                             'Meiryo UI', 24, center=True, color_rgb=(150, 150, 150))
        else:
            add_text_box(slide, f"No Image: {label}", SLIDE_WIDTH_PX/2, SLIDE_HEIGHT_PX/2,
                         'Meiryo UI', 24, center=True, color_rgb=(150, 150, 150))

    # ä¿å­˜
    output_path = os.path.join(temp_dir, f"order_{order_data['order_id']}.pptx")
    prs.save(output_path)
    print(f"[Canva] PowerPoint created: {output_path}")
    return output_path


def create_pdf(order_data, temp_dir):
    """PDFã‚’ä½œæˆï¼ˆé€æ˜åº¦å¯¾å¿œç‰ˆï¼‰"""
    print(f"[Canva] Creating PDF for order #{order_data['order_id']}...")

    sim_data = order_data['sim_data']
    sim_image = order_data['sim_image']
    groom = sim_data.get('groomName', '')
    bride = sim_data.get('brideName', '')

    # PDFãƒšãƒ¼ã‚¸ã‚µã‚¤ã‚ºï¼ˆãƒã‚¤ãƒ³ãƒˆå˜ä½ã€1000x1000pxç›¸å½“ï¼‰
    PAGE_SIZE = (1000, 1000)

    output_path = os.path.join(temp_dir, f"order_{order_data['order_id']}.pdf")
    c = pdf_canvas.Canvas(output_path, pagesize=PAGE_SIZE)

    # cutoutç”»åƒã®URLã‚’å–å¾—
    background = sim_data.get('background', 'product')
    cutout_urls = {
        'product': find_cutout_url(order_data['board_name'], order_data['board_number'], order_data['board_size'], 'product'),
        'product_back': find_cutout_url(order_data['board_name'], order_data['board_number'], order_data['board_size'], 'product_back'),
        'noclear': find_cutout_url(order_data['board_name'], order_data['board_number'], order_data['board_size'], 'noclear'),
        'noclear_back': find_cutout_url(order_data['board_name'], order_data['board_number'], order_data['board_size'], 'noclear_back'),
    }

    base_font = sim_data.get('baseFont', 'Alex Brush')
    # ãƒ†ã‚­ã‚¹ãƒˆè‰²: 'burn'=èŒ¶è‰², 'white'=ç™½
    text_color_name = sim_data.get('textColor', 'burn')
    text_color_hex = (1.0, 1.0, 1.0) if text_color_name == 'white' else (42/255, 24/255, 16/255)

    # ========== 1ãƒšãƒ¼ã‚¸ç›®: ã‚·ãƒŸãƒ¥ãƒ¬ãƒ¼ã‚·ãƒ§ãƒ³ç”»åƒ + æ³¨æ–‡æƒ…å ± ==========
    if sim_image:
        try:
            if sim_image.startswith('http'):
                response = requests.get(sim_image, timeout=30)
                img = Image.open(BytesIO(response.content))
            else:
                if ',' in sim_image:
                    sim_image = sim_image.split(',')[1]
                img_data = base64.b64decode(sim_image)
                img = Image.open(BytesIO(img_data))

            # ä¸€æ™‚ä¿å­˜
            sim_path = os.path.join(temp_dir, 'sim_image.png')
            img.save(sim_path, 'PNG')

            # PPTXã¨åŒã˜ãƒ­ã‚¸ãƒƒã‚¯ã§å¤§ããé…ç½®
            img_width, img_height = img.size
            available_height = PAGE_SIZE[1] * 0.80  # 80%ã®é«˜ã•ã‚’ä½¿ç”¨
            available_width = PAGE_SIZE[0] * 0.95   # 95%ã®å¹…ã‚’ä½¿ç”¨
            scale = min(available_width / img_width, available_height / img_height)
            draw_width = img_width * scale
            draw_height = img_height * scale

            # ä¸­å¤®ä¸Šå¯„ã›é…ç½®ï¼ˆPDFã¯å·¦ä¸‹åŸç‚¹ï¼‰
            x = (PAGE_SIZE[0] - draw_width) / 2
            y = PAGE_SIZE[1] - draw_height - 20  # ä¸Šã‹ã‚‰20pxä¸‹

            c.drawImage(sim_path, x, y, width=draw_width, height=draw_height)
        except Exception as e:
            print(f"[WARN] Sim image error: {e}")

    # æ³¨æ–‡æƒ…å ±ãƒ†ã‚­ã‚¹ãƒˆ
    c.setFillColorRGB(0.16, 0.16, 0.16)
    c.setFont("Helvetica", 14)
    info_y = 150
    c.drawCentredString(PAGE_SIZE[0]/2, info_y, f"Order #{order_data['order_id']} - {order_data['board_name']} No.{order_data['board_number']}")
    c.setFont("Helvetica", 11)
    c.drawCentredString(PAGE_SIZE[0]/2, info_y - 25, f"Groom: {groom}  Bride: {bride}")
    c.drawCentredString(PAGE_SIZE[0]/2, info_y - 50, f"Date: {order_data['wedding_date']}")
    # ãƒ•ã‚©ãƒ³ãƒˆè¡¨ç¤ºåï¼ˆã‚·ãƒŸãƒ¥ãƒ¬ãƒ¼ã‚¿ãƒ¼ã¨åŒã˜ã€Œã€œé¢¨ã€è¡¨è¨˜ã‚’ä½¿ç”¨ï¼‰
    font_display_name = FONT_DISPLAY_MAP.get(base_font, base_font)
    c.drawCentredString(PAGE_SIZE[0]/2, info_y - 75, f"Font: {font_display_name}")

    c.showPage()

    # ========== 2ãƒšãƒ¼ã‚¸ç›®: èƒŒæ™¯cutout + ãƒ†ã‚­ã‚¹ãƒˆ + ãƒ„ãƒªãƒ¼ ==========
    # PPTXã¨åŒã˜ã‚¹ã‚±ãƒ¼ãƒ«ä¿‚æ•°ã‚’ä½¿ç”¨
    FONT_SCALE = PAGE_SIZE[0] / 500  # = 2.0
    board_size_pct = sim_data.get('boardSize', 130) / 100

    # æ¿ã®å¢ƒç•Œæƒ…å ±ï¼ˆãƒ‡ãƒ•ã‚©ãƒ«ãƒˆå€¤ã€ç”»åƒè§£æå¾Œã«æ›´æ–°ï¼‰
    board_bounds = {'minX': 0, 'maxX': 1, 'minY': 0, 'maxY': 1}
    board_center_offset_x = 0
    board_center_offset_y = 0
    draw_w = 0
    draw_h = 0

    # èƒŒæ™¯cutoutç”»åƒ
    cutout_url = cutout_urls.get(background, cutout_urls['product'])
    try:
        response = requests.get(cutout_url, timeout=30)
        if response.status_code == 200:
            cutout_img = Image.open(BytesIO(response.content))
            if cutout_img.mode != 'RGBA':
                cutout_img = cutout_img.convert('RGBA')

            # æ¿ã®å®Ÿéš›ã®å¢ƒç•Œã‚’æ¤œå‡ºï¼ˆé€ééƒ¨åˆ†ã‚’é™¤ãï¼‰
            img_w, img_h = cutout_img.size
            alpha = cutout_img.split()[3]
            alpha_data = alpha.load()

            # éé€æ˜ãƒ”ã‚¯ã‚»ãƒ«ã®å¢ƒç•Œã‚’æ¤œå‡º
            min_x, max_x, min_y, max_y = img_w, 0, img_h, 0
            for py in range(img_h):
                for px in range(img_w):
                    if alpha_data[px, py] > 10:  # é–¾å€¤10ä»¥ä¸Šã‚’æ¿ã¨ã—ã¦èªè­˜
                        min_x = min(min_x, px)
                        max_x = max(max_x, px)
                        min_y = min(min_y, py)
                        max_y = max(max_y, py)

            # æ­£è¦åŒ–ï¼ˆ0-1ã®ç¯„å›²ã«ï¼‰
            if max_x > min_x and max_y > min_y:
                board_bounds = {
                    'minX': min_x / img_w,
                    'maxX': max_x / img_w,
                    'minY': min_y / img_h,
                    'maxY': max_y / img_h
                }
                print(f"[PDF] Board bounds: X={board_bounds['minX']:.2f}-{board_bounds['maxX']:.2f}, Y={board_bounds['minY']:.2f}-{board_bounds['maxY']:.2f}")

            cutout_path = os.path.join(temp_dir, 'cutout_bg.png')
            cutout_img.save(cutout_path, 'PNG')

            # ã‚·ãƒŸãƒ¥ãƒ¬ãƒ¼ã‚¿ãƒ¼ã¨åŒã˜ãƒ­ã‚¸ãƒƒã‚¯ã§æç”»ã‚µã‚¤ã‚ºã‚’è¨ˆç®—
            base_scale = 0.95
            img_ratio = img_w / img_h
            max_width = PAGE_SIZE[0] * base_scale
            max_height = PAGE_SIZE[1] * base_scale

            # ã‚¢ã‚¹ãƒšã‚¯ãƒˆæ¯”ã‚’ç¶­æŒã—ã¦ãƒ•ã‚£ãƒƒãƒˆ
            if img_w / max_width > img_h / max_height:
                base_width = max_width
                base_height = max_width / img_ratio
            else:
                base_height = max_height
                base_width = max_height * img_ratio

            # boardSizeã‚’é©ç”¨
            draw_w = base_width * board_size_pct
            draw_h = base_height * board_size_pct
            img_x = (PAGE_SIZE[0] - draw_w) / 2
            img_y = (PAGE_SIZE[1] - draw_h) / 2

            # å®Ÿéš›ã®æ¿ã®ä¸­å¿ƒã‚ªãƒ•ã‚»ãƒƒãƒˆã‚’è¨ˆç®—ï¼ˆç”»åƒä¸­å¿ƒã‹ã‚‰ã®ãšã‚Œï¼‰
            board_center_offset_x = ((board_bounds['minX'] + board_bounds['maxX']) / 2 - 0.5) * draw_w
            board_center_offset_y = ((board_bounds['minY'] + board_bounds['maxY']) / 2 - 0.5) * draw_h
            print(f"[PDF] Board center offset: ({board_center_offset_x:.1f}, {board_center_offset_y:.1f})")

            c.drawImage(cutout_path, img_x, img_y, width=draw_w, height=draw_h, mask='auto')
            print(f"[PDF] Cutout: draw={draw_w:.0f}x{draw_h:.0f}, boardSize={board_size_pct*100:.0f}%")
    except Exception as e:
        print(f"[WARN] Cutout error: {e}")

    # å®Ÿéš›ã®æ¿ã®ã‚µã‚¤ã‚ºï¼ˆé€ééƒ¨åˆ†ã‚’é™¤ãï¼‰
    actual_board_w = draw_w * (board_bounds['maxX'] - board_bounds['minX'])
    actual_board_h = draw_h * (board_bounds['maxY'] - board_bounds['minY'])

    # æ¿ã®ä¸­å¿ƒä½ç½®ï¼ˆãƒšãƒ¼ã‚¸ä¸Šã§ã®å®Ÿéš›ã®ä½ç½®ï¼‰
    board_center_x = PAGE_SIZE[0] / 2 + board_center_offset_x
    board_center_y = PAGE_SIZE[1] / 2 - board_center_offset_y  # Yè»¸åè»¢

    # ãƒ†ã‚­ã‚¹ãƒˆè¦ç´ ï¼ˆæ¿ã®å®Ÿéš›ã®å¢ƒç•Œã‚’åŸºæº–ã«é…ç½®ï¼‰
    c.setFillColorRGB(*text_color_hex)

    # ã‚¿ã‚¤ãƒˆãƒ« - æ¿ã®å¢ƒç•Œã‚’åŸºæº–ã«é…ç½®
    title_key = sim_data.get('title', 'wedding')
    title_text = sim_data.get('customTitle', '') if title_key == 'custom' else TITLES.get(title_key, 'Wedding Certificate')
    # titleY=22 ã¯æ¿ã®ä¸Šã‹ã‚‰22%ã®ä½ç½®
    title_x = board_center_x + actual_board_w * ((sim_data.get('titleX', 50) - 50) / 100)
    title_y_pct = sim_data.get('titleY', 22) / 100
    title_y = board_center_y + actual_board_h / 2 - actual_board_h * title_y_pct
    title_size = 24 * (sim_data.get('titleSize', 100) / 100) * FONT_SCALE
    c.setFont("Helvetica", title_size)
    c.drawCentredString(title_x, title_y, title_text)

    # æœ¬æ–‡
    template_key = sim_data.get('template', 'holy')
    body_text = sim_data.get('customText', '') if template_key == 'custom' else TEMPLATES.get(template_key, '')
    body_x = board_center_x + actual_board_w * ((sim_data.get('bodyX', 50) - 50) / 100)
    body_y_pct = sim_data.get('bodyY', 32) / 100
    body_y = board_center_y + actual_board_h / 2 - actual_board_h * body_y_pct
    body_size = 11 * (sim_data.get('bodySize', 115) / 100) * FONT_SCALE
    body_line_height = sim_data.get('bodyLineHeight', 1.4)
    c.setFont("Helvetica", body_size)
    for i, line in enumerate(body_text.split('\n')):
        c.drawCentredString(body_x, body_y - i * body_size * body_line_height, line.strip())

    # æ—¥ä»˜
    date_format_key = sim_data.get('dateFormat', 'western')
    formatted_date = sim_data.get('customDate', '') if date_format_key == 'custom' else format_date(order_data['wedding_date'], date_format_key)
    date_x = board_center_x + actual_board_w * ((sim_data.get('dateX', 50) - 50) / 100)
    date_y_pct = sim_data.get('dateY', 60) / 100
    date_y = board_center_y + actual_board_h / 2 - actual_board_h * date_y_pct
    date_size = 18 * (sim_data.get('dateSize', 85) / 100) * FONT_SCALE
    c.setFont("Helvetica", date_size)
    c.drawCentredString(date_x, date_y, formatted_date)

    # åå‰
    name_x = board_center_x + actual_board_w * ((sim_data.get('nameX', 50) - 50) / 100)
    name_y_pct = sim_data.get('nameY', 74) / 100
    name_y = board_center_y + actual_board_h / 2 - actual_board_h * name_y_pct
    name_size = 32 * (sim_data.get('nameSize', 90) / 100) * FONT_SCALE
    c.setFont("Helvetica", name_size)
    name_text = f"{groom}  &  {bride}"
    c.drawCentredString(name_x, name_y, name_text)

    # ãƒ„ãƒªãƒ¼ç”»åƒï¼ˆé€æ˜åº¦ä¿æŒï¼‰- æ¿ã‚µã‚¤ã‚ºã‚’åŸºæº–ã«è¨ˆç®—
    if sim_data.get('showTree', False):
        tree_type = sim_data.get('treeType', 'simple')
        tree_x_pct = sim_data.get('treeX', 0.75)
        tree_y_pct = sim_data.get('treeY', 0.65)
        tree_size_pct = sim_data.get('treeSize', 80) / 100

        tree_url = f"{TREE_IMAGES_URL}/tree-{tree_type}.png"
        print(f"[PDF] Downloading tree: {tree_url}")
        print(f"[PDF] Tree params: x={tree_x_pct}, y={tree_y_pct}, size={tree_size_pct*100}%")
        print(f"[PDF] Board size: {actual_board_w:.0f}x{actual_board_h:.0f}")

        try:
            response = requests.get(tree_url, timeout=30)
            if response.status_code == 200:
                tree_img = Image.open(BytesIO(response.content))
                print(f"[PDF] Tree original: {tree_img.size}, mode={tree_img.mode}")

                # RGBAç¢ºä¿
                if tree_img.mode != 'RGBA':
                    tree_img = tree_img.convert('RGBA')

                # ã‚¢ãƒ«ãƒ•ã‚¡ãƒã‚§ãƒƒã‚¯
                alpha = tree_img.split()[3]
                print(f"[PDF] Tree alpha range: {alpha.getextrema()}")

                tree_path = os.path.join(temp_dir, 'tree.png')
                tree_img.save(tree_path, 'PNG')

                # ===== æ¿ã‚µã‚¤ã‚ºåŸºæº–ã®ãƒ„ãƒªãƒ¼ã‚µã‚¤ã‚ºè¨ˆç®— =====
                # è¦–è¦šçš„èª¿æ•´: 0.35=å°, 0.7=ä¸­, 1.0=é©æ­£, 1.2=å¤§
                tree_base_ratio = 1.0  # è¦–è¦šçš„ã«ç¢ºèªæ¸ˆã¿
                draw_w = actual_board_w * tree_size_pct * tree_base_ratio

                # ã‚¢ã‚¹ãƒšã‚¯ãƒˆæ¯”ã‚’ç¶­æŒ
                tree_w, tree_h = tree_img.size
                aspect_ratio = tree_h / tree_w
                draw_h = draw_w * aspect_ratio

                print(f"[PDF] Tree draw size: {draw_w:.0f}x{draw_h:.0f} ({draw_w/actual_board_w*100:.1f}% of board)")

                # ===== æ¿ã‚’åŸºæº–ã«ã—ãŸä½ç½®è¨ˆç®— =====
                # treeX, treeY ã¯æ¿å†…ã§ã®ç›¸å¯¾ä½ç½®ï¼ˆ0-1ï¼‰
                # æ¿ã®å·¦ä¸Šã‚’(0,0)ã€å³ä¸‹ã‚’(1,1)ã¨ã™ã‚‹
                # æ¿ã®å®Ÿéš›ã®å·¦ç«¯ãƒ»ä¸Šç«¯ã‚’è¨ˆç®—
                board_left = board_center_x - actual_board_w / 2
                board_top = board_center_y + actual_board_h / 2  # PDFåº§æ¨™ç³»ï¼ˆYè»¸ä¸Šå‘ãï¼‰

                # ãƒ„ãƒªãƒ¼ä¸­å¿ƒä½ç½®ï¼ˆæ¿å†…åº§æ¨™ï¼‰
                tree_center_x = board_left + actual_board_w * tree_x_pct
                tree_center_y = board_top - actual_board_h * tree_y_pct  # Yè»¸åè»¢

                # å·¦ä¸‹åº§æ¨™ã«å¤‰æ›
                x = tree_center_x - draw_w / 2
                y = tree_center_y - draw_h / 2

                print(f"[PDF] Tree position: center=({tree_center_x:.0f}, {tree_center_y:.0f}), corner=({x:.0f}, {y:.0f})")

                # mask='auto' ã§PNGé€æ˜åº¦ã‚’è‡ªå‹•é©ç”¨
                c.drawImage(tree_path, x, y, width=draw_w, height=draw_h, mask='auto')
                print(f"[PDF] Tree: size={tree_size_pct*100:.0f}%, draw={draw_w:.0f}x{draw_h:.0f}, pos=({x:.0f}, {y:.0f})")
        except Exception as e:
            print(f"[WARN] Tree error: {e}")
            import traceback
            traceback.print_exc()

    c.showPage()

    # ========== 3-6ãƒšãƒ¼ã‚¸ç›®: cutoutç”»åƒ ==========
    cutout_labels = [
        ('product', 'Front Clear'),
        ('product_back', 'Back Clear'),
        ('noclear', 'Front No-Clear'),
        ('noclear_back', 'Back No-Clear'),
    ]

    for key, label in cutout_labels:
        try:
            response = requests.get(cutout_urls[key], timeout=30)
            if response.status_code == 200:
                img = Image.open(BytesIO(response.content))
                if img.mode != 'RGBA':
                    img = img.convert('RGBA')

                max_size = 800
                if max(img.size) > max_size:
                    ratio = max_size / max(img.size)
                    new_size = (int(img.size[0] * ratio), int(img.size[1] * ratio))
                    img = img.resize(new_size, Image.LANCZOS)

                img_path = os.path.join(temp_dir, f'{key}.png')
                img.save(img_path, 'PNG')

                img_w, img_h = img.size
                max_draw = PAGE_SIZE[0] * 0.9
                scale = min(max_draw / img_w, max_draw / img_h)
                draw_w = img_w * scale
                draw_h = img_h * scale
                x = (PAGE_SIZE[0] - draw_w) / 2
                y = (PAGE_SIZE[1] - draw_h) / 2

                c.drawImage(img_path, x, y, width=draw_w, height=draw_h, mask='auto')
        except Exception as e:
            print(f"[WARN] {label} error: {e}")
            c.setFont("Helvetica", 20)
            c.drawCentredString(PAGE_SIZE[0]/2, PAGE_SIZE[1]/2, f"No Image: {label}")

        c.showPage()

    c.save()
    print(f"[Canva] PDF created: {output_path}")
    return output_path


def send_discord_notification(order_data, design, webhook_url, order=None):
    """Discordé€šçŸ¥é€ä¿¡ï¼ˆæ–°è¦æ³¨æ–‡ + Canvaãƒªãƒ³ã‚¯çµ±åˆç‰ˆï¼‰"""
    if not webhook_url:
        return False

    edit_url = design.get("urls", {}).get("edit_url", "")
    groom = order_data['sim_data'].get('groomName', '')
    bride = order_data['sim_data'].get('brideName', '')

    # æ³¨æ–‡æƒ…å ±ã‚’å–å¾—ï¼ˆorderã‚ªãƒ–ã‚¸ã‚§ã‚¯ãƒˆãŒã‚ã‚‹å ´åˆï¼‰
    customer_name = ""
    order_total = ""
    payment_method = ""
    customer_phone = ""
    customer_email = ""

    if order:
        billing = order.get('billing', {})
        customer_name = f"{billing.get('last_name', '')} {billing.get('first_name', '')}"
        order_total = order.get('total', '0')
        payment_method = order.get('payment_method_title', '')
        customer_phone = billing.get('phone', '')
        customer_email = billing.get('email', '')

    embed = {
        "title": f"ğŸ›’ æ–°è¦æ³¨æ–‡ #{order_data['order_id']}",
        "color": 0x06C755,  # LINEç·‘
        "fields": [
            {"name": "ğŸ‘¤ ãŠå®¢æ§˜", "value": customer_name or f"{groom} & {bride}", "inline": True},
            {"name": "ğŸ’° é‡‘é¡", "value": f"Â¥{int(float(order_total)):,}" if order_total else "N/A", "inline": True},
            {"name": "ğŸ’³ æ”¯æ‰•æ–¹æ³•", "value": payment_method or "N/A", "inline": True},
            {"name": "ğŸ“¦ å•†å“", "value": f"{order_data['board_name']} No.{order_data['board_number']}", "inline": False},
            {"name": "ğŸ“… æŒ™å¼æ—¥", "value": order_data['wedding_date'], "inline": False},
            {"name": "ğŸ“ é€£çµ¡å…ˆ", "value": f"TEL: {customer_phone}\nEmail: {customer_email}" if customer_phone else "N/A", "inline": False},
            {"name": "ğŸ¨ Canva", "value": f"[ãƒ‡ã‚¶ã‚¤ãƒ³ã‚’ç·¨é›†ã™ã‚‹]({edit_url})", "inline": False},
        ],
        "footer": {"text": "i.tategu è‡ªå‹•åŒ–ã‚·ã‚¹ãƒ†ãƒ ï¼ˆRailwayï¼‰"},
    }

    # å•†å“ç”»åƒãŒã‚ã‚Œã°ã‚µãƒ ãƒã‚¤ãƒ«ã«è¨­å®š
    if order:
        for item in order.get('line_items', []):
            image_url = item.get('image', {}).get('src', '')
            if image_url:
                embed['thumbnail'] = {'url': image_url}
                break

    response = requests.post(
        webhook_url,
        json={"embeds": [embed]},
        headers={"Content-Type": "application/json"}
    )

    return response.status_code == 204


# ç™ºé€ç®¡ç†ãƒãƒ£ãƒ³ãƒãƒ«ID
DISCORD_SHIPPING_CHANNEL_ID = "1463452139312644240"

def send_shipping_notification(order_data, order, bot_token):
    """ç™ºé€ç®¡ç†ãƒãƒ£ãƒ³ãƒãƒ«ã«ä½æ‰€æƒ…å ±ã‚’æŠ•ç¨¿"""
    if not order or not bot_token:
        print("[Shipping] Missing order or bot_token")
        return False

    billing = order.get('billing', {})
    shipping = order.get('shipping', {})

    # ç™ºé€å…ˆæƒ…å ±ï¼ˆshippingå„ªå…ˆã€ãªã‘ã‚Œã°billingï¼‰
    postcode = shipping.get('postcode') or billing.get('postcode', '')
    state = shipping.get('state') or billing.get('state', '')
    city = shipping.get('city') or billing.get('city', '')
    address1 = shipping.get('address_1') or billing.get('address_1', '')
    address2 = shipping.get('address_2') or billing.get('address_2', '')

    # éƒ½é“åºœçœŒã‚³ãƒ¼ãƒ‰å¤‰æ›ï¼ˆç°¡æ˜“ç‰ˆï¼‰
    JP_STATES = {
        'JP01': 'åŒ—æµ·é“', 'JP02': 'é’æ£®çœŒ', 'JP03': 'å²©æ‰‹çœŒ', 'JP04': 'å®®åŸçœŒ',
        'JP05': 'ç§‹ç”°çœŒ', 'JP06': 'å±±å½¢çœŒ', 'JP07': 'ç¦å³¶çœŒ', 'JP08': 'èŒ¨åŸçœŒ',
        'JP09': 'æ ƒæœ¨çœŒ', 'JP10': 'ç¾¤é¦¬çœŒ', 'JP11': 'åŸ¼ç‰çœŒ', 'JP12': 'åƒè‘‰çœŒ',
        'JP13': 'æ±äº¬éƒ½', 'JP14': 'ç¥å¥ˆå·çœŒ', 'JP15': 'æ–°æ½ŸçœŒ', 'JP16': 'å¯Œå±±çœŒ',
        'JP17': 'çŸ³å·çœŒ', 'JP18': 'ç¦äº•çœŒ', 'JP19': 'å±±æ¢¨çœŒ', 'JP20': 'é•·é‡çœŒ',
        'JP21': 'å²é˜œçœŒ', 'JP22': 'é™å²¡çœŒ', 'JP23': 'æ„›çŸ¥çœŒ', 'JP24': 'ä¸‰é‡çœŒ',
        'JP25': 'æ»‹è³€çœŒ', 'JP26': 'äº¬éƒ½åºœ', 'JP27': 'å¤§é˜ªåºœ', 'JP28': 'å…µåº«çœŒ',
        'JP29': 'å¥ˆè‰¯çœŒ', 'JP30': 'å’Œæ­Œå±±çœŒ', 'JP31': 'é³¥å–çœŒ', 'JP32': 'å³¶æ ¹çœŒ',
        'JP33': 'å²¡å±±çœŒ', 'JP34': 'åºƒå³¶çœŒ', 'JP35': 'å±±å£çœŒ', 'JP36': 'å¾³å³¶çœŒ',
        'JP37': 'é¦™å·çœŒ', 'JP38': 'æ„›åª›çœŒ', 'JP39': 'é«˜çŸ¥çœŒ', 'JP40': 'ç¦å²¡çœŒ',
        'JP41': 'ä½è³€çœŒ', 'JP42': 'é•·å´çœŒ', 'JP43': 'ç†Šæœ¬çœŒ', 'JP44': 'å¤§åˆ†çœŒ',
        'JP45': 'å®®å´çœŒ', 'JP46': 'é¹¿å…å³¶çœŒ', 'JP47': 'æ²–ç¸„çœŒ'
    }
    state_name = JP_STATES.get(state, state)

    full_address = f"{state_name}{city}{address1}"
    if address2:
        full_address += f" {address2}"

    customer_name = f"{billing.get('last_name', '')} {billing.get('first_name', '')}"
    customer_phone = billing.get('phone', '')
    order_total = order.get('total', '0')
    payment_method = order.get('payment_method_title', '')

    # å•†å“å
    products = []
    for item in order.get('line_items', []):
        products.append(item.get('name', ''))
    product_names = ', '.join(products) if products else order_data.get('board_name', '')

    embed = {
        "title": f"ğŸŸ¡ æœªç™ºé€ | #{order_data['order_id']} {customer_name} æ§˜",
        "color": 0xFFD700,  # é»„è‰²
        "fields": [
            {"name": "ğŸ“ é›»è©±", "value": customer_phone or "N/A", "inline": True},
            {"name": "ğŸ“¦ å•†å“", "value": product_names, "inline": True},
            {"name": "ğŸ’° é‡‘é¡", "value": f"Â¥{int(float(order_total)):,} / {payment_method}", "inline": True},
            {"name": "ã€’ ä½æ‰€", "value": f"{postcode} {full_address}" if postcode else full_address, "inline": False},
        ],
    }

    # Discord Bot APIã§é€ä¿¡
    url = f"https://discord.com/api/v10/channels/{DISCORD_SHIPPING_CHANNEL_ID}/messages"
    headers = {
        "Authorization": f"Bot {bot_token}",
        "Content-Type": "application/json"
    }

    try:
        response = requests.post(url, json={"embeds": [embed]}, headers=headers)
        if response.status_code in [200, 201]:
            print(f"[Shipping] Notification sent for order #{order_data['order_id']}")
            return True
        else:
            print(f"[Shipping] Failed: {response.status_code} - {response.text}")
            return False
    except Exception as e:
        print(f"[Shipping] Error: {e}")
        return False


def clear_processing_lock(order_id, wc_url, wc_key, wc_secret):
    """å‡¦ç†ä¸­ãƒ­ãƒƒã‚¯ã‚’è§£é™¤ï¼ˆå¤±æ•—æ™‚ç”¨ï¼‰"""
    url = f"{wc_url}/wp-json/wc/v3/orders/{order_id}?consumer_key={wc_key}&consumer_secret={wc_secret}"
    try:
        requests.put(url, json={"meta_data": [{"key": "canva_processing", "value": ""}]})
        print(f"[Canva] Lock released for order #{order_id}")
    except Exception as e:
        print(f"[WARN] Failed to release lock: {e}")


def send_discord_error_notification(order_id, error_message, webhook_url):
    """Discord ã‚¨ãƒ©ãƒ¼é€šçŸ¥é€ä¿¡"""
    if not webhook_url:
        return False

    embed = {
        "title": f"âš ï¸ Canvaå‡¦ç†ã‚¨ãƒ©ãƒ¼ #{order_id}",
        "color": 15158332,  # èµ¤è‰²
        "fields": [
            {"name": "ã‚¨ãƒ©ãƒ¼å†…å®¹", "value": str(error_message)[:500]},
        ],
        "footer": {"text": "i.tategu Canvaè‡ªå‹•åŒ–ï¼ˆRailwayï¼‰"},
    }

    try:
        response = requests.post(
            webhook_url,
            json={"embeds": [embed]},
            headers={"Content-Type": "application/json"}
        )
        return response.status_code == 204
    except:
        return False


def mark_order_processed(order_id, design_url, wc_url, wc_key, wc_secret):
    """æ³¨æ–‡ã‚’å‡¦ç†æ¸ˆã¿ã«ãƒãƒ¼ã‚¯ + ã‚¹ãƒ†ãƒ¼ã‚¿ã‚¹ã‚’ã€Œãƒ‡ã‚¶ã‚¤ãƒ³æ‰“ã¡åˆã‚ã›ä¸­ã€ã«å¤‰æ›´"""
    # WooCommerce REST APIã¯ã‚¯ã‚¨ãƒªãƒ‘ãƒ©ãƒ¡ãƒ¼ã‚¿ã§èªè¨¼
    url = f"{wc_url}/wp-json/wc/v3/orders/{order_id}?consumer_key={wc_key}&consumer_secret={wc_secret}"

    data = {
        "status": "designing",  # ãƒ‡ã‚¶ã‚¤ãƒ³æ‰“ã¡åˆã‚ã›ä¸­ (short slug for WP 20-char limit)
        "meta_data": [
            {"key": "canva_automation_done", "value": "1"},
            {"key": "canva_processing", "value": ""},  # ãƒ­ãƒƒã‚¯è§£é™¤
            {"key": "canva_design_url", "value": design_url},
        ]
    }

    try:
        response = requests.put(url, json=data)
        print(f"[WC Update] Status: {response.status_code}")
        if response.status_code != 200:
            print(f"[WC Update] Error: {response.text[:500]}")
            return False
        print(f"[WC Update] Order #{order_id} marked as processed, status â†’ designing")
        return True
    except Exception as e:
        print(f"[WC Update] Exception: {e}")
        return False


def process_order(order_id, config):
    """
    æ³¨æ–‡ã‚’å‡¦ç†ã—ã¦Canvaãƒ‡ã‚¶ã‚¤ãƒ³ã‚’ä½œæˆ

    config: {
        'wc_url': WooCommerce URL,
        'wc_key': Consumer Key,
        'wc_secret': Consumer Secret,
        'canva_access_token': Canva Access Token,
        'canva_refresh_token': Canva Refresh Token,
        'discord_webhook': Discord Webhook URL,
    }
    """
    print(f"\n{'='*50}")
    print(f"[Canva] Processing order #{order_id}")
    print(f"{'='*50}")

    # æ³¨æ–‡å–å¾—
    order = get_order_from_woocommerce(order_id, config['wc_url'], config['wc_key'], config['wc_secret'])
    if not order:
        print(f"[ERROR] Order not found: {order_id}")
        return False

    # æ—¢ã«å‡¦ç†æ¸ˆã¿ or å‡¦ç†ä¸­ã‹ãƒã‚§ãƒƒã‚¯
    meta = {m['key']: m['value'] for m in order.get('meta_data', [])}
    if meta.get('canva_automation_done') or meta.get('canva_processing'):
        print(f"[SKIP] Already processed or in progress: {order_id}")
        return False

    # å³åº§ã«å‡¦ç†ä¸­ãƒ•ãƒ©ã‚°ã‚’ç«‹ã¦ã‚‹ï¼ˆé‡è¤‡é˜²æ­¢ãƒ­ãƒƒã‚¯ï¼‰
    lock_acquired = False
    try:
        lock_url = f"{config['wc_url']}/wp-json/wc/v3/orders/{order_id}?consumer_key={config['wc_key']}&consumer_secret={config['wc_secret']}"
        requests.put(lock_url, json={"meta_data": [{"key": "canva_processing", "value": "1"}]})
        lock_acquired = True
        print(f"[Canva] Lock acquired for order #{order_id}")
    except Exception as e:
        print(f"[WARN] Lock failed: {e}")

    # ãƒ­ãƒƒã‚¯å–å¾—å¾Œã®å‡¦ç†ï¼ˆå¤±æ•—æ™‚ã¯å¿…ãšãƒ­ãƒƒã‚¯è§£é™¤ï¼‰
    success = False
    error_message = None

    try:
        # æ³¨æ–‡ãƒ‡ãƒ¼ã‚¿è§£æ
        order_data = parse_order_data(order)

        if not order_data['board_name']:
            error_message = "No board info"
            print(f"[SKIP] No board info: {order_id}")
            return False

        print(f"[Canva] Product: {order_data['board_name']} No.{order_data['board_number']}")
        print(f"[Canva] Names: {order_data['sim_data'].get('groomName', '')} & {order_data['sim_data'].get('brideName', '')}")

        # ä¸€æ™‚ãƒ‡ã‚£ãƒ¬ã‚¯ãƒˆãƒªä½œæˆ
        with tempfile.TemporaryDirectory() as temp_dir:
            # PDFä½œæˆï¼ˆé€æ˜åº¦å¯¾å¿œç‰ˆï¼‰
            pdf_path = create_pdf(order_data, temp_dir)

            # Canvaã‚¿ã‚¤ãƒˆãƒ«
            groom = order_data['sim_data'].get('groomName', '')
            bride = order_data['sim_data'].get('brideName', '')
            canva_title = f"æ³¨æ–‡{order_id} {order_data['board_name']} No.{order_data['board_number']} {groom}ï¼†{bride} {order_data['wedding_date']}"

            # Canvaã‚¤ãƒ³ãƒãƒ¼ãƒˆ
            print(f"[Canva] Importing PDF to Canva...")
            design, error_info = import_to_canva(
                pdf_path, canva_title,
                config['canva_access_token'],
                config['canva_refresh_token']
            )

            if not design:
                error_message = f"Canva import failed: {error_info}"
                print(f"[ERROR] {error_message}")
                return False

            design_id = design.get('id')
            print(f"[Canva] Design ID: {design_id}")

            # Discordé€šçŸ¥ï¼ˆæ³¨æ–‡æƒ…å ±+Canvaãƒªãƒ³ã‚¯çµ±åˆç‰ˆï¼‰
            print(f"[Canva] Sending Discord notification...")
            send_discord_notification(order_data, design, config['discord_webhook'], order)
            print(f"[Canva] Discord notification sent")

            # ç™ºé€ç®¡ç†ãƒãƒ£ãƒ³ãƒãƒ«ã¸ä½æ‰€æƒ…å ±é€šçŸ¥
            bot_token = config.get('discord_bot_token', '')
            if bot_token:
                print(f"[Canva] Sending shipping notification...")
                send_shipping_notification(order_data, order, bot_token)
            else:
                print(f"[WARN] No bot token, skipping shipping notification")

            # å‡¦ç†æ¸ˆã¿ãƒãƒ¼ã‚¯ï¼ˆã“ã“ã§ãƒ­ãƒƒã‚¯ã‚‚è§£é™¤ã•ã‚Œã‚‹ï¼‰
            design_url = design.get('urls', {}).get('edit_url', '')
            mark_order_processed(order_id, design_url, config['wc_url'], config['wc_key'], config['wc_secret'])
            print(f"[Canva] Order marked as processed")
            success = True

        print(f"[Canva] Order #{order_id} completed!")
        return True

    except Exception as e:
        error_message = str(e)
        print(f"[ERROR] Processing failed: {e}")
        import traceback
        traceback.print_exc()
        return False

    finally:
        # å¤±æ•—æ™‚ã¯ãƒ­ãƒƒã‚¯è§£é™¤ & ã‚¨ãƒ©ãƒ¼é€šçŸ¥
        if lock_acquired and not success:
            clear_processing_lock(order_id, config['wc_url'], config['wc_key'], config['wc_secret'])
            if error_message and config.get('discord_webhook'):
                send_discord_error_notification(order_id, error_message, config['discord_webhook'])
